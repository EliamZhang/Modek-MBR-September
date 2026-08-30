# -*- coding: utf-8 -*-
"""在 Model MBR August.pptx 中新增「澳洲模型体系规划」页(插到第 2 位)。"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

PPTX = r"c:\Users\zhangyuliang02\Desktop\临时文件\9月MBR\Model MBR August.pptx"

# 色彩
INK = RGBColor(0x0A, 0x0A, 0x0A)      # 标题黑
BODY = RGBColor(0x3D, 0x3C, 0x3A)     # 正文
GRAY = RGBColor(0x60, 0x60, 0x60)     # 次要
NAVY = RGBColor(0x1F, 0x4E, 0x79)     # 深蓝
GOLD = RGBColor(0xD6, 0xA0, 0x00)     # 金
LBLUE = RGBColor(0xE3, 0xEB, 0xF4)    # 浅蓝
LGOLD = RGBColor(0xFA, 0xF3, 0xE0)    # 浅金
LGRAY = RGBColor(0xF0, 0xF0, 0xF0)    # 浅灰
LINEC = RGBColor(0xD6, 0xD6, 0xD6)    # 弱化线
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"


def set_run(run, size, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_para(tf, runs, first=False, align=PP_ALIGN.LEFT, space_after=0,
             line_spacing=1.1, anchor=None):
    """runs: list of (text, size, bold, color)"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        set_run(r, size, bold, color)
    return p


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def add_rect(slide, x, y, w, h, fill, line=None, line_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_topbar(slide, x, y, w):
    return add_rect(slide, x, y, w, 0.06, NAVY)


def add_section_title(slide, y, text):
    add_rect(slide, 0.56, y + 0.02, 0.07, 0.22, GOLD)
    box, tf = add_textbox(slide, 0.70, y, 9.0, 0.30)
    add_para(tf, [(text, 12, True, INK)], first=True, line_spacing=1.0)


def main():
    prs = Presentation(PPTX)

    # 新页用空白 layout(与第 1 页相同),避免占位符干扰
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # ---- 页面标题 / 引导语 / 分隔线 ----
    box, tf = add_textbox(slide, 0.66, 0.23, 9.72, 0.40)
    add_para(tf, [("体系规划:四项目标,一个底座,三大领域", 20, True, INK)],
             first=True, line_spacing=1.3)

    box, tf = add_textbox(slide, 0.56, 0.64, 11.92, 0.48)
    add_para(tf, [("规模、结构、质量、效率四项目标并重;数据基建打底,投放、运营、风险三大领域共享同一分层体系。", 18, False, GRAY)],
             first=True, line_spacing=1.4)

    add_rect(slide, 0.56, 1.14, 11.92, 0.012, LINEC)

    # ---- 章节标题 1:四项目标 ----
    add_section_title(slide, 1.28, "四项目标:规模 · 结构 · 质量 · 效率")

    # ---- 4 个目标卡 ----
    goals = [
        ("规模", "支撑放款增长", "拓展 APP、新渠道与复贷"),
        ("结构", "NON SAAC ≥ 40%", "非 SACC 产品占比目标"),
        ("质量", "风险 ≤ 7.1%", "整体风险水平控制"),
        ("效率", "LTV 统筹经营", "统筹获客、额度、定价、权益"),
    ]
    gw, gh, gap, gy = 2.85, 1.08, 0.19, 1.60
    for i, (name, metric, desc) in enumerate(goals):
        x = 0.56 + i * (gw + gap)
        card = add_rect(slide, x, gy, gw, gh, LBLUE)
        add_topbar(slide, x, gy, gw)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.16)
        tf.margin_bottom = Inches(0.08)
        add_para(tf, [(name, 10.5, True, INK)], first=True, line_spacing=1.0, space_after=4)
        add_para(tf, [(metric, 13, True, NAVY)], line_spacing=1.0, space_after=4)
        add_para(tf, [(desc, 8.5, False, BODY)], line_spacing=1.0)

    # ---- 底座:数据基建 ----
    by, bh = 2.92, 1.10
    base = add_rect(slide, 0.56, by, 11.92, bh, LGRAY)
    add_topbar(slide, 0.56, by, 11.92)
    tf = base.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.06)
    add_para(tf, [("数据基建 · 一个底座", 11, True, NAVY),
                  ("  自主流水识别 + 统一数据底座,为模型、策略与归因提供一致口径", 9.5, False, BODY)],
             first=True, line_spacing=1.0, space_after=2)
    for item in ["BS-CAT 生产化与并行验证",
                 "打通 IDP2 全链路:IDA → IDP → BS-CAT → S-CALC → 落库 → 风险决策 → 人工复核",
                 "商户 / 类别知识库智能体"]:
        add_para(tf, [("★ ", 9, True, GOLD), (item, 9, True, NAVY)],
                 line_spacing=1.05, space_after=1)

    # ---- 章节标题 2:三大领域 ----
    add_section_title(slide, 4.10, "三大领域:共享同一分层体系")

    # ---- 3 个领域卡 ----
    fields = [
        ("投放市场", [
            ("人群圈选:风险 × 意愿 × 价值分层输出人群包", False),
            ("Lead Partner 线索评估", False),
            ("投放归因基建", True),
            ("回传监控基建", True),
            ("边际 CPS 与扩量决策", True),
            ("Campaign 早期质量预测", True),
        ]),
        ("用户运营", [
            ("盈利性评估", False),
            ("自动寻优定价", False),
            ("Fundo Score", False),
            ("触达归因与效果监控", True),
            ("触达策略与人群编排", True),
            ("运营监控 AI 智能体", True),
            ("结清客户复贷召回", True),
            ("沉睡客户识别唤醒", True),
            ("竞品经营洞察 AI 智能体", True),
        ]),
        ("风险管理", [
            ("新老客 PD 模型迭代:增强头尾区分,平衡通过率、客群结构与风险", True),
            ("风险 × 价值联合额度策略:多模型交叉分箱自动寻优,输出最优额度提升客群与产品承接矩阵标签", True),
        ]),
    ]
    fw, fh, fgap, fy = 3.87, 2.23, 0.16, 4.42
    for i, (title, items) in enumerate(fields):
        x = 0.56 + i * (fw + fgap)
        card = add_rect(slide, x, fy, fw, fh, WHITE, line=LINEC, line_w=1.5)
        add_topbar(slide, x, fy, fw)
        add_rect(slide, x + 0.16, fy + 0.18, 0.06, 0.20, GOLD)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.28)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.14)
        tf.margin_bottom = Inches(0.10)
        add_para(tf, [(title, 12, True, INK)], first=True, line_spacing=1.0, space_after=6)
        for text, star in items:
            if star:
                add_para(tf, [("★ ", 9, True, GOLD), (text, 9, True, NAVY)],
                         line_spacing=1.1, space_after=3)
            else:
                add_para(tf, [(text, 9, False, BODY)], line_spacing=1.1, space_after=3)

    # ---- 结论条 ----
    cy, ch = 6.82, 0.42
    concl = add_rect(slide, 0.56, cy, 11.92, ch, LGOLD)
    add_rect(slide, 0.56, cy, 0.07, ch, GOLD)
    tf = concl.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.24)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    add_para(tf, [("四项目标统领全局——数据基建打底,投放、运营、风险三大领域在统一数据与分层体系上协同建设。", 9, True, INK)],
             first=True, line_spacing=1.1)

    # 页序:新页(在最后)移到第 2 位(slide 1 之后)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    last = ids[-1]
    sldIdLst.remove(last)
    ids[1].addnext(last)

    prs.save(PPTX)
    print("saved:", PPTX)


if __name__ == "__main__":
    main()
