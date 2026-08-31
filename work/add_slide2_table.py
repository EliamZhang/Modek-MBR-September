# -*- coding: utf-8 -*-
"""向第2页(价值模型P1)左下半幅空位补充"表3:平均可支配盈余";并将全 deck 表格编号顺移。
表3(第2页,原为来源材料表4) / 原表3(第3页,风险×价值)->表4 / 原表4(第5页,Top10)->表5。
样式与 build_report_pages.py 保持一致(极简横线表, 表头灰底, 数据白底)。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

PPTX = r"c:\Users\zhangyuliang02\Desktop\临时文件\9月MBR\Model MBR August.pptx"

INK = RGBColor(0x0A, 0x0A, 0x0A)
BODY = RGBColor(0x3D, 0x3C, 0x3A)
GRAY = RGBColor(0x60, 0x60, 0x60)
LGRAY = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"


def set_run(run, size, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_para(tf, runs, first=False, align=PP_ALIGN.LEFT, space_after=0, line_spacing=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        set_run(r, size, bold, color)
    return p


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def _cell_borders(cell, header=False):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)

    def mk(tag, w, color=None):
        ln = etree.Element(qn(tag))
        ln.set('w', str(w))
        ln.set('cap', 'flat')
        if color is None:
            etree.SubElement(ln, qn('a:noFill'))
        else:
            sf = etree.SubElement(ln, qn('a:solidFill'))
            c = etree.SubElement(sf, qn('a:srgbClr'))
            c.set('val', color)
        return ln

    tcPr.insert(0, mk('a:lnL', 12700))
    tcPr.insert(1, mk('a:lnR', 12700))
    tcPr.insert(2, mk('a:lnT', 12700))
    if header:
        tcPr.insert(3, mk('a:lnB', 19050, '0A0A0A'))
    else:
        tcPr.insert(3, mk('a:lnB', 9525, 'D6D6D6'))


def add_table(slide, x, y, w, data, col_widths, row_h, header_size=8, data_size=8,
              center_cols=None):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(row_h * rows))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0')
    tblPr.set('bandRow', '0')
    for el in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(el)
    for i, cw_ in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw_)
    for r in range(rows):
        tbl.rows[r].height = Inches(row_h)
    center_cols = center_cols or set()
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = LGRAY if r == 0 else WHITE
            _cell_borders(cell, header=(r == 0))
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c in center_cols else PP_ALIGN.LEFT
            p.line_spacing = 1.0
            txt = data[r][c]
            run = p.add_run()
            run.text = txt if txt else ' '
            set_run(run, header_size if r == 0 else data_size, bold=(r == 0),
                    color=INK if r == 0 else BODY)


def renumber_caption(slide, old_prefix, new_prefix):
    """把形状标题文本框中的 old_prefix(如 '表3：')改为 new_prefix。"""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith(old_prefix):
            tf = sh.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.text.startswith(old_prefix):
                        run.text = new_prefix + run.text[len(old_prefix):]
            return True
    return False


def main():
    prs = Presentation(PPTX)
    s = prs.slides[1]  # 第2页 = 价值模型P1

    # 表3(盈余-申请样本) 标题: 与表1标题同高对齐
    box, tf = add_textbox(s, 0.56, 3.86, 5.90, 0.25)
    add_para(tf, [("表3：各申请金额区间 × 价值分箱的平均可支配盈余", 10.5, True, INK)],
             first=True, line_spacing=1.0)

    d3 = [
        ["申请金额区间", "价值1", "价值2", "价值3", "价值4", "价值5"],
        ["500–1000", "1,106", "763", "288", "793", "287"],
        ["1000–1500", "940", "685", "731", "603", "46"],
        ["1500–2000", "1,096", "798", "729", "608", "144"],
        ["2000–2500", "1,708", "737", "496", "828", "−154"],
        ["2500–5000", "1,196", "583", "681", "569", "124"],
        ["5000+", "1,488", "931", "989", "1,010", "−353"],
        ["合计", "1,334", "745", "598", "761", "84"],
    ]
    add_table(s, 0.56, 4.14, 4.84, d3, [1.30, 0.708, 0.708, 0.708, 0.708, 0.708],
              0.25, center_cols={1, 2, 3, 4, 5})

    # 注: 首尾判读 + 价值5转负
    box, tf = add_textbox(s, 0.56, 6.22, 5.84, 0.30)
    add_para(tf, [("注：盈余＝收入−支出（可支配余量）；按「价值1最高、价值5最低」首尾判读，中间档并非严格单调；价值5在2000–2500与5000+两段转负（入不敷出）。",
                   7.5, False, GRAY)], first=True, line_spacing=1.0)

    # 编号顺移: 第3页 表3->表4, 第5页 表4->表5
    ok3 = renumber_caption(prs.slides[2], "表3：", "表4：")
    ok5 = renumber_caption(prs.slides[4], "表4：", "表5：")
    prs.save(PPTX)
    print("saved")
    print("slide3 表3->表4:", ok3)
    print("slide5 表4->表5:", ok5)


if __name__ == "__main__":
    main()
