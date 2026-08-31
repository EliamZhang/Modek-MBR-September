# -*- coding: utf-8 -*-
"""修正第2页标题/注文本框宽度:
- 新"表3"标题 5.90 -> 4.84(收进左列, 不与"表1"标题相撞)
- 既有"表1"标题 11.92 -> 6.20(收进右半幅, 消除越界; 历史遗留尺寸)
- "表3"注 宽 5.84->4.84, 高 0.30->0.40(容纳两行, 避免压到结论条)
"""
from pptx import Presentation
from pptx.util import Inches

PPTX = r"c:\Users\zhangyuliang02\Desktop\临时文件\9月MBR\Model MBR August.pptx"


def resize_by_text(slide, keyword, width, height=None):
    for sh in slide.shapes:
        if sh.has_text_frame and keyword in sh.text_frame.text:
            sh.width = Inches(width)
            if height is not None:
                sh.height = Inches(height)
            return True, sh.text_frame.text.strip()[:20]
    return False, None


def main():
    prs = Presentation(PPTX)
    s = prs.slides[1]
    a = resize_by_text(s, "表3：", 4.84)
    b = resize_by_text(s, "表1：各申请金额", 6.20)
    c = resize_by_text(s, "注：盈余", 4.84, 0.40)
    prs.save(PPTX)
    print("表3标题:", a)
    print("表1标题:", b)
    print("表3注:", c)


if __name__ == "__main__":
    main()
