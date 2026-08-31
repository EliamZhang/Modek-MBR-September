# -*- coding: utf-8 -*-
"""按《02-价值模型汇报-PPT正式版.md》P4-P7 同步月会PPT的画像内容:
1) 第5页(旧画像页)整页重建为新 P4 画像对比页(表5);
2) 第5页后新增 P5 高价值但高风险(表6)、P6 低价值但低风险(表7)、附录 Top10(表8)三页;
3) 旧页 Top10 表(表5)内容迁移到附录页,并按 MD 加回「变量」列(旧 PPT 曾省略该列);
   旧注、旧结论按新 MD 替换。
表编号: 原表5 Top10 → 表8; 新表5/6/7 = 画像对比表。"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout.reconfigure(encoding='utf-8')

PPTX = r"c:\Users\zhangyuliang02\Desktop\临时文件\9月MBR\月会PPT.pptx"

INK = RGBColor(0x0A, 0x0A, 0x0A)
BODY = RGBColor(0x3D, 0x3C, 0x3A)
GRAY = RGBColor(0x60, 0x60, 0x60)
NAVY = RGBColor(0x1F, 0x4E, 0x79)
GOLD = RGBColor(0xD6, 0xA0, 0x00)
LGOLD = RGBColor(0xFA, 0xF3, 0xE0)
LGRAY = RGBColor(0xF0, 0xF0, 0xF0)
LINEC = RGBColor(0xD6, 0xD6, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"


def set_run(run, size, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_para(tf, runs, first=False, align=PP_ALIGN.LEFT, space_after=0, line_spacing=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        set_run(r, size, bold, color)
    return p


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def add_rect(slide, x, y, w, h, fill, line=None, line_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def page_head(slide, title, guide):
    box, tf = add_textbox(slide, 0.66, 0.23, 11.92, 0.40)
    add_para(tf, [(title, 20, True, INK)], first=True, line_spacing=1.3)
    box, tf = add_textbox(slide, 0.56, 0.64, 11.92, 0.48)
    add_para(tf, [(guide, 18, False, GRAY)], first=True, line_spacing=1.4)
    add_rect(slide, 0.56, 1.14, 11.92, 0.012, LINEC)


def page_concl(slide, text):
    cy, ch = 6.82, 0.42
    concl = add_rect(slide, 0.56, cy, 11.92, ch, LGOLD)
    add_rect(slide, 0.56, cy, 0.07, ch, GOLD)
    tf = concl.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.24)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    add_para(tf, [(text, 9, True, INK)], first=True, line_spacing=1.1)


def cw_est(ch, pt):
    if ord(ch) > 0x2E7F:
        return pt / 72.0
    if 0x20 <= ord(ch) <= 0x7E:
        return pt / 72.0 * 0.52
    return pt / 72.0 * 0.9


def add_pill(slide, x, y, w, h, text, size=9):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = LGOLD
    pill.line.color.rgb = GOLD
    pill.line.width = Pt(1)
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = text
    set_run(r, size, True, NAVY)
    return pill


def tags_row(slide, texts, y=1.30, h=0.28, gap=0.18, size=9):
    x = 0.56
    for t in texts:
        w = sum(cw_est(c, size) for c in t) + 0.28
        add_pill(slide, x, y, w, h, t, size=size)
        x += w + gap


def _cell_borders(cell, header=False):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)

    def mk(tag, w, color=None):
        ln = etree.Element(qn(tag))
        ln.set('w', str(w))
        ln.set('cap', 'flat')
        if color is None:
            etree.SubElement(ln, qn('a:noFill'))
        else:
            sf = etree.SubElement(ln, qn('a:solidFill'))
            c = etree.SubElement(sf, qn('a:srgbClr'))
            c.set('val', color)
        return ln

    tcPr.insert(0, mk('a:lnL', 12700))
    tcPr.insert(1, mk('a:lnR', 12700))
    tcPr.insert(2, mk('a:lnT', 12700))
    if header:
        tcPr.insert(3, mk('a:lnB', 19050, '0A0A0A'))
    else:
        tcPr.insert(3, mk('a:lnB', 9525, 'D6D6D6'))


def add_table(slide, x, y, w, data, col_widths, row_heights, header_size=8, data_size=8,
              center_cols=None, num_cols=None, dim_col=None):
    """row_heights: 每行高度(含表头)。num_cols: 深蓝加粗的数值列;dim_col: 灰色弱化列(解读)。"""
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w),
                                Inches(sum(row_heights)))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0')
    tblPr.set('bandRow', '0')
    for el in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(el)
    for i, cw_ in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw_)
    for r in range(rows):
        tbl.rows[r].height = Inches(row_heights[r])
    center_cols = center_cols or set()
    num_cols = num_cols or set()
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = LGRAY if r == 0 else WHITE
            _cell_borders(cell, header=(r == 0))
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c in center_cols else PP_ALIGN.LEFT
            p.line_spacing = 1.0
            txt = data[r][c]
            run = p.add_run()
            run.text = txt if txt else ' '
            if r == 0:
                set_run(run, header_size, True, INK)
            elif c in num_cols:
                set_run(run, data_size, True, NAVY)
            elif c == dim_col:
                set_run(run, data_size, False, GRAY)
            else:
                set_run(run, data_size, False, BODY)


def clear_slide(slide):
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)


def build_page(slide, page):
    page_head(slide, page["title"], page["guide"])
    if page["tags"]:
        tags_row(slide, page["tags"])
    box, tf = add_textbox(slide, 0.56, 1.74, 11.92, 0.25)
    add_para(tf, [(page["caption"], 10.5, True, INK)], first=True, line_spacing=1.0)
    data = [page["header"]] + page["rows"]
    add_table(slide, 0.56, 2.04, 11.92, data, page["col_widths"], page["row_heights"],
              center_cols=page["center_cols"], num_cols=page["num_cols"],
              dim_col=page.get("dim_col"))
    if page.get("note"):
        box, tf = add_textbox(slide, 0.56, 4.80, 11.92, 0.22)
        add_para(tf, [(page["note"], 8, False, GRAY)], first=True, line_spacing=1.0)
    page_concl(slide, page["concl"])


def add_page(prs, layout, page):
    s = prs.slides.add_slide(layout)
    build_page(s, page)
    return s


P4 = dict(
    title="价值画像由收入体量、工资连续性与职业可识别性共同决定",
    guide="价值差更像「职业现金流弱」，并非负债更重或多头更多。",
    tags=["高收入体量", "工资更连续", "主业更清晰", "福利收入依赖低"],
    caption="表5：画像对比（价值 1/2 档 vs 4/5 档）",
    header=["画像维度", "价值好", "价值差", "画像解读"],
    rows=[
        ["工资收入", "41,139", "11,130", "价值好约为 3.7 倍"],
        ["总收入", "140,807", "59,771", "整体资金体量更高"],
        ["工资入账次数", "20.5", "10.6", "工资现金流更连续"],
        ["主行业入账", "40,916", "16,067", "职业主收入更强"],
        ["Centrelink 收入占比", "2.1%", "14.6%", "价值差更依赖福利收入"],
        ["发薪周期标准差", "2.72", "3.79", "价值差发薪更不稳定"],
        ["行业识别缺失率", "3.5%", "30.7%", "价值差职业信息更弱"],
        ["可支配盈余", "1,334", "84", "价值差可支配空间有限"],
    ],
    col_widths=[2.00, 1.86, 1.86, 6.20],
    row_heights=[0.20] + [0.30] * 8,
    center_cols={1, 2},
    num_cols={1, 2},
    dim_col=3,
    note="结构补充：价值好人群 couple 占比更高（43.8% vs 29.1%），Lead Partner / Owned Channels 占比也更高；该差异仅作结构描述，不解释为因果。",
    concl="信贷还款金额与次数在价值差人群中反而更低（5,963 / 13.5 vs 20,439 / 29.4），当前数据不支持「价值差 = 高负债 / 多头更多」。",
)

P5 = dict(
    title="高价值但高风险：收入更高，但资金波动与风险消费同步放大",
    guide="这类客户的价值来自收入端，风险则来自高支出、高波动和账户异常。",
    tags=["高收入", "高支出", "高赌博", "高波动"],
    caption="表6：画像对比（价值好且风险好 vs 价值好但风险差）",
    header=["画像维度", "价值好且风险好", "价值好但风险差", "风险信号"],
    rows=[
        ["工资收入", "34,024", "50,253", "收入能力并不弱"],
        ["总收入", "96,943", "191,383", "资金体量约 2 倍"],
        ["总支出", "92,348", "193,711", "收支几乎打平"],
        ["赌博支出", "2,398 / 2.5%", "7,416 / 5.2%", "金额与占比同时升高"],
        ["余额波动标准差", "4,947", "8,442", "账户波动显著放大"],
        ["余额异常天数", "3.82", "5.13", "账户异常更多"],
        ["收入波动 CV", "1.83", "2.10", "收入稳定性更弱"],
    ],
    col_widths=[2.00, 1.86, 1.86, 6.20],
    row_heights=[0.20] + [0.34] * 7,
    center_cols={1, 2},
    num_cols={1, 2},
    dim_col=3,
    note="注：赌博支出为金额 / 占比。",
    concl="不宜因价值高直接扩额；优先用小额、短期产品承接，并持续观察赌博支出、余额异常和收入波动。",
)

P6 = dict(
    title="低价值但低风险：收入偏弱，但支出克制、账户稳定、还款更活跃",
    guide="这类客户不是高收入客群，而是行为更保守、风险相对可控的低收入客群。",
    tags=["低收入", "低支出", "低赌博", "低波动", "有还款活动"],
    caption="表7：画像对比（价值差但风险好 vs 价值差且风险差）",
    header=["画像维度", "价值差但风险好", "价值差且风险差", "风险信号"],
    rows=[
        ["工资收入", "8,750", "13,314", "低风险并非来自高收入"],
        ["Centrelink 收入占比", "22.1%", "5.9%", "收入结构偏福利类"],
        ["总支出", "43,062", "75,496", "支出更克制"],
        ["赌博支出", "531.7 / 1.3%", "2,067 / 3.7%", "风险消费更低"],
        ["余额波动标准差", "478.1", "1,181", "账户波动更小"],
        ["余额异常天数", "3.45", "5.22", "账户异常更少"],
        ["信贷还款金额", "7,582", "4,133", "还款活动更活跃"],
        ["贷款还款次数", "15.9", "10.8", "还款活动相对更多"],
    ],
    col_widths=[2.00, 1.86, 1.86, 6.20],
    row_heights=[0.20] + [0.30] * 8,
    center_cols={1, 2},
    num_cols={1, 2},
    dim_col=3,
    note="注：赌博支出为金额 / 占比。",
    concl="不宜因价值弱一刀切拒绝；可用低额度 Personal Loan 验证其额度使用、复贷与真实贡献。",
)

APP = dict(
    title="模型由收入、现金流、负债与支出等多类变量共同驱动",
    guide="Top 3 均为收入与现金流特征，重要度合计 51.08%。",
    tags=[],
    caption="表8：入模变量重要度 Top10",
    header=["排名", "变量", "含义", "重要度"],
    rows=[
        ["1", "bank_txn_income_wages_sum_182d", "近 182 天工资收入汇总金额", "20.25%"],
        ["2", "bank_txn_income_global_mean_56d", "近 56 天收入均值", "15.60%"],
        ["3", "bank_txn_balance_debit_mean_168d", "近 168 天借方金额均值", "15.23%"],
        ["4", "bank_txn_lender_repay_bnpl_amount_sum", "先买后付（BNPL）还款金额总和", "6.98%"],
        ["5", "bank_txn_category_cluster_transfers_share_28d", "近 28 天转账簇金额占比", "6.32%"],
        ["6", "bank_txn_lender_repay_loan_count_max", "个人贷款还款单机构最大累计笔数", "4.98%"],
        ["7", "bank_txn_balance_cv_56d", "近 56 天余额变异系数", "4.91%"],
        ["8", "bank_txn_lender_disburse_loan_amount_avg_l168d", "近 168 天个人贷款放款金额均值", "4.75%"],
        ["9", "bank_txn_category_insurance_debit_amt_14d", "近 14 天保险出账金额", "4.56%"],
        ["10", "bank_txn_category_cluster_debt_avg_ticket_168d", "债务支出簇近 168 天平均客单价", "3.60%"],
    ],
    col_widths=[0.90, 3.70, 6.02, 1.30],
    row_heights=[0.20] + [0.24] * 10,
    center_cols={0, 3},
    num_cols={3},
    dim_col=None,
    note=None,
    concl="画像用于解释分层，不替代风险判断；地域、渠道、具体行业与余额绝对值仅作结构补充，不直接形成强干预结论。",
)


def main():
    prs = Presentation(PPTX)
    assert len(prs.slides) == 7, f"预期 7 页, 实际 {len(prs.slides)}"

    # ---- 0. 旧第5页 Top10 表与 MD 核对(数据不可删改) ----
    old = prs.slides[4]
    old_rows = None
    for sh in old.shapes:
        if sh.has_table:
            tbl = sh.table
            old_rows = [[tbl.cell(r, c).text for c in range(3)] for r in range(len(tbl.rows))]
    md_rows = [[r[0], r[2], r[3]] for r in APP["rows"]]
    same = old_rows is not None and old_rows[0] == APP["header"][:3] and old_rows[1:] == md_rows
    print("旧 Top10 表与 MD 一致:", same)
    if not same:
        print("旧表:", old_rows)
        print("MD :", [APP["header"][:3]] + md_rows)

    # ---- 1. 第5页重建为新 P4 ----
    clear_slide(old)
    build_page(old, P4)
    print("第5页已重建为新 P4(表5)")

    # ---- 2. 新增 P5 / P6 / 附录 三页 ----
    layout = next(l for l in prs.slide_layouts if l.name == "自定义版式")
    new_slides = [add_page(prs, layout, P5), add_page(prs, layout, P6), add_page(prs, layout, APP)]
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    anchor = ids[4]  # 第5页
    for el in ids[-3:]:
        sldIdLst.remove(el)
        anchor.addnext(el)
        anchor = el
    print("已新增 P5 / P6 / 附录 三页并插到第5页之后")

    prs.save(PPTX)
    print("saved:", PPTX)
    prs2 = Presentation(PPTX)
    for i, s in enumerate(prs2.slides, 1):
        title = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                title = sh.text_frame.text.strip().split("\n")[0]
                break
        print(f"  slide {i}: {title}")


if __name__ == "__main__":
    main()
