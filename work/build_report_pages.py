# -*- coding: utf-8 -*-
"""向 Model MBR August.pptx 追加 6 页:价值模型 P1-P4 + 银行流水 P1-P2(深蓝+金)。
布局按内容估算高度排版,避免重叠与越界。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

PPTX = r"c:\Users\zhangyuliang02\Desktop\临时文件\9月MBR\Model MBR August.pptx"

INK = RGBColor(0x0A, 0x0A, 0x0A)
BODY = RGBColor(0x3D, 0x3C, 0x3A)
GRAY = RGBColor(0x60, 0x60, 0x60)
NAVY = RGBColor(0x1F, 0x4E, 0x79)
GOLD = RGBColor(0xD6, 0xA0, 0x00)
LBLUE = RGBColor(0xE3, 0xEB, 0xF4)
LGOLD = RGBColor(0xFA, 0xF3, 0xE0)
LGRAY = RGBColor(0xF0, 0xF0, 0xF0)
LINEC = RGBColor(0xD6, 0xD6, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"


def set_run(run, size, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_para(tf, runs, first=False, align=PP_ALIGN.LEFT, space_after=0, line_spacing=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        set_run(r, size, bold, color)
    return p


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def add_rect(slide, x, y, w, h, fill, line=None, line_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_topbar(slide, x, y, w):
    return add_rect(slide, x, y, w, 0.06, NAVY)


def page_head(slide, title, guide):
    box, tf = add_textbox(slide, 0.66, 0.23, 11.92, 0.40)
    add_para(tf, [(title, 20, True, INK)], first=True, line_spacing=1.3)
    box, tf = add_textbox(slide, 0.56, 0.64, 11.92, 0.48)
    add_para(tf, [(guide, 18, False, GRAY)], first=True, line_spacing=1.4)
    add_rect(slide, 0.56, 1.14, 11.92, 0.012, LINEC)


def page_concl(slide, text):
    cy, ch = 6.82, 0.42
    concl = add_rect(slide, 0.56, cy, 11.92, ch, LGOLD)
    add_rect(slide, 0.56, cy, 0.07, ch, GOLD)
    tf = concl.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.24)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    add_para(tf, [(text, 9, True, INK)], first=True, line_spacing=1.1)


def cw_est(ch, pt):
    """字宽估算(与验证脚本同口径)。"""
    if ord(ch) > 0x2E7F:
        return pt / 72.0
    if 0x20 <= ord(ch) <= 0x7E:
        return pt / 72.0 * 0.52
    return pt / 72.0 * 0.9


def est_lines(text, pt, avail):
    if not text.strip():
        return 0
    cur = 0.0
    n = 1
    for ch in text:
        cur += cw_est(ch, pt)
        if cur > avail:
            n += 1
            cur = cw_est(ch, pt)
    return n


def block(slide, x, y, w, title, body, body_size=9, title_size=10):
    """金色小竖条 + 标题 + 正文;按估算行数返回占高,文本框按内容设高。"""
    avail_w = w - 0.14
    lines = 0
    for runs in body:
        text = ''.join(r[0] for r in runs)
        lines += est_lines(text, body_size, avail_w)
    h = 0.20 + lines * (body_size / 72.0 * 1.15) + 0.10  # 标题行+正文+缓冲
    add_rect(slide, x, y + 0.03, 0.06, 0.18, GOLD)
    box, tf = add_textbox(slide, x + 0.14, y, avail_w, h)
    add_para(tf, [(title, title_size, True, INK)], first=True, line_spacing=1.0, space_after=2)
    for runs in body:
        add_para(tf, runs, line_spacing=1.05, space_after=2)
    return h


def _cell_borders(cell, header=False):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)

    def mk(tag, w, color=None):
        ln = etree.Element(qn(tag))
        ln.set('w', str(w))
        ln.set('cap', 'flat')
        if color is None:
            etree.SubElement(ln, qn('a:noFill'))
        else:
            sf = etree.SubElement(ln, qn('a:solidFill'))
            c = etree.SubElement(sf, qn('a:srgbClr'))
            c.set('val', color)
        return ln

    tcPr.insert(0, mk('a:lnL', 12700))
    tcPr.insert(1, mk('a:lnR', 12700))
    tcPr.insert(2, mk('a:lnT', 12700))
    if header:
        tcPr.insert(3, mk('a:lnB', 19050, '0A0A0A'))
    else:
        tcPr.insert(3, mk('a:lnB', 9525, 'D6D6D6'))


def add_table(slide, x, y, w, data, col_widths, row_h, header_size=8, data_size=8,
              center_cols=None):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(row_h * rows))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0')
    tblPr.set('bandRow', '0')
    for el in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(el)
    for i, cw_ in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw_)
    for r in range(rows):
        tbl.rows[r].height = Inches(row_h)
    center_cols = center_cols or set()
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = LGRAY if r == 0 else WHITE
            _cell_borders(cell, header=(r == 0))
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c in center_cols else PP_ALIGN.LEFT
            p.line_spacing = 1.0
            txt = data[r][c]
            run = p.add_run()
            run.text = txt if txt else ' '
            set_run(run, header_size if r == 0 else data_size, bold=(r == 0),
                    color=INK if r == 0 else BODY)


def stat_card(slide, x, y, w, h, label, value, sub=None):
    card = add_rect(slide, x, y, w, h, LBLUE)
    add_topbar(slide, x, y, w)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.06)
    add_para(tf, [(label, 9, True, INK)], first=True, line_spacing=1.0, space_after=3)
    add_para(tf, [(value, 15, True, NAVY)], line_spacing=1.0, space_after=2)
    if sub:
        add_para(tf, [(sub, 8, False, GRAY)], line_spacing=1.0)


def main():
    prs = Presentation(PPTX)
    blank = prs.slide_layouts[0]

    # ================= 价值模型 P1 =================
    s = prs.slides.add_slide(blank)
    page_head(s, "模型识别的是客户价值，不是历史额度策略",
              "控制额度后，模型仍能稳定区分客户表现——价值信号存在于额度之外。")
    box, tf = add_textbox(s, 0.56, 1.30, 4.84, 2.45)
    add_para(tf, [("直接回答", 10, True, INK)], first=True, line_spacing=1.0, space_after=2)
    add_para(tf, [("额度在先、利息收入在后；但同等额度能否转化为收入，取决于客户的收入、现金流与履约能力——模型识别的是客户本身的价值潜力。", 9, False, BODY)], line_spacing=1.05, space_after=6)
    add_para(tf, [("为什么必须检验", 10, True, INK)], line_spacing=1.0, space_after=2)
    add_para(tf, [("利息收入 ≈ 在贷余额 × 利率 × 用款时长，与额度机械关联；「高额度 = 高价值」的质疑仅凭论证无法排除。", 9, False, BODY)], line_spacing=1.05, space_after=6)
    add_para(tf, [("决定性检验——同额度分层", 10, True, INK)], line_spacing=1.0, space_after=2)
    add_para(tf, [("收入 6 个申请金额区间全部单调（5000+ 档：", 9, False, BODY),
                  ("10,388 vs 3,891", 9, True, NAVY),
                  ("）；逾期率 4 个总额度区间首尾方向全部随价值变差上升（", 9, False, BODY),
                  ("4.37% vs 20.59%", 9, True, NAVY),
                  ("）；盈余 1,334→84、支出 7,902→3,566 同样单调。", 9, False, BODY)], line_spacing=1.05, space_after=6)
    add_para(tf, [("局限披露", 10, True, INK)], line_spacing=1.0, space_after=2)
    add_para(tf, [("总额度 2,500+ 区间暂无足够放款表现样本（高额度申请基本被拒），逾期维度暂不能验证；申请金额维度在该区间仍保持收入单调，不影响结论方向。", 9, False, BODY)], line_spacing=1.05)
    add_rect(s, 5.55, 1.30, 0.01, 4.4, LINEC)
    box, tf = add_textbox(s, 6.28, 1.30, 6.42, 0.25)
    add_para(tf, [("表2：各总额度区间 × 价值分箱的 3M30 逾期率", 10.5, True, INK)], first=True, line_spacing=1.0)
    d2 = [
        ["总额度区间", "价值1", "价值2", "价值3", "价值4", "价值5", "有效样本数"],
        ["500–1000", "4.37%", "8.04%", "7.47%", "12.22%", "20.59%", "5,769"],
        ["1000–1500", "5.73%", "8.35%", "10.94%", "20.14%", "28.04%", "3,748"],
        ["1500–2000", "7.80%", "10.76%", "13.55%", "16.84%", "37.88%", "1,254"],
        ["2000–2500", "8.84%", "12.08%", "17.17%", "24.71%", "28.12%", "2,181"],
        ["合计", "7.05%", "9.46%", "10.61%", "15.91%", "22.90%", "12,988"],
    ]
    add_table(s, 6.28, 1.58, 6.42, d2, [1.30, 0.80, 0.80, 0.80, 0.80, 0.80, 1.12],
              0.24, center_cols={1, 2, 3, 4, 5, 6})
    box, tf = add_textbox(s, 6.28, 3.10, 6.42, 0.4)
    add_para(tf, [("注：500–1000 段价值2/3 有轻微倒挂（8.04% vs 7.47%），上屏引用「首尾方向全部正确」，不写「全部严格单调」。", 7.5, False, GRAY)],
             first=True, line_spacing=1.0)
    box, tf = add_textbox(s, 0.56, 3.86, 11.92, 0.25)
    add_para(tf, [("表1：各申请金额区间 × 价值分箱的平均收入", 10.5, True, INK)], first=True, line_spacing=1.0)
    d1 = [
        ["申请金额区间", "价值1", "价值2", "价值3", "价值4", "价值5"],
        ["500–1000", "9,559", "6,926", "6,081", "5,518", "3,857"],
        ["1000–1500", "8,596", "6,850", "6,322", "5,446", "3,982"],
        ["1500–2000", "8,829", "7,103", "6,186", "5,535", "3,928"],
        ["2000–2500", "9,891", "7,258", "6,465", "5,826", "4,028"],
        ["2500–5000", "9,288", "7,454", "6,374", "5,467", "3,965"],
        ["5000+", "10,388", "7,849", "6,768", "6,322", "3,891"],
        ["合计", "9,611", "7,284", "6,366", "5,682", "3,924"],
    ]
    add_table(s, 0.56, 4.14, 11.92, d1, [1.52, 2.08, 2.08, 2.08, 2.08, 2.08],
              0.25, center_cols={1, 2, 3, 4, 5})
    page_concl(s, "额度决定价值释放的上限，价值决定额度投向哪里；「额度贡献了多少价值」属因果问题，历史数据无法拆分，留待增量额度试点实验回答（方案另行汇报）。")

    # ================= 价值模型 P2 =================
    s = prs.slides.add_slide(blank)
    page_head(s, "效果：跨时间稳定，分箱单调，交叉分化清晰",
              "风险与价值两个维度各自排序成立，交叉后两端分化近 5 倍——联合分层具备数据基础。")
    stat_card(s, 0.56, 1.32, 2.88, 0.98, "AUC（OOT）", "0.7294", "INS 0.7015")
    stat_card(s, 3.56, 1.32, 2.88, 0.98, "KS（OOT）", "0.3467", "INS 0.2939")
    stat_card(s, 6.56, 1.32, 2.88, 0.98, "Top10% 弱价值率", "63.00%", "Lift 2.45")
    stat_card(s, 9.56, 1.32, 2.88, 0.98, "Bottom10%", "7.30%", "高分 vs 低分")
    box, tf = add_textbox(s, 0.56, 2.48, 11.92, 1.55)
    add_para(tf, [("模型概况：", 9, True, INK),
                  ("新客价值模型 1.0，LightGBM 二分类，17 个银行流水特征，建模样本 56,330（成交客户，2024-04 起）；目标=「3 个月利息收入低于 160」的弱价值标记（占比 27.12%）；原始分越高=弱价值概率越高，对外价值 1（最好）→ 5（最弱）。", 9, False, BODY)],
             first=True, line_spacing=1.05, space_after=5)
    add_para(tf, [("分箱单调：", 9, True, INK),
                  ("申请样本 151,579，价值 1→5：收入 9,611→3,924、通过率 41.04%→6.57%、3M30 逾期率 7.05%→22.90%，逐档单调、无倒挂。", 9, False, BODY)],
             line_spacing=1.05, space_after=5)
    add_para(tf, [("风险模型对照：", 9, True, INK),
                  ("风险 1→5，3M30 逾期率 5.92%→22.02%、通过率 32.91%→7.49%，排序清晰、无倒挂——两个维度各自成立，联合分层成立。", 9, False, BODY)],
             line_spacing=1.05, space_after=5)
    add_para(tf, [("交叉结果与使用边界：", 9, True, INK),
                  ("申请样本 157,741，最好格「风险1×价值1」5.74%、最差格「风险5×价值5」27.05%；样本集中于最差格（46,235，29.3%），价值5档占 46.45%。预测性模型回答「当前策略下谁更可能贡献价值」，OOT 弱价值概率低估约 2.15pp，适合排序分层，概率阈值/定价需先校准。", 9, False, BODY)],
             line_spacing=1.05)
    box, tf = add_textbox(s, 0.56, 4.12, 11.92, 0.25)
    add_para(tf, [("表3：风险 × 价值 3M30 人数逾期率", 10.5, True, INK)], first=True, line_spacing=1.0)
    d3 = [
        ["价值 \\ 风险", "风险1", "风险2", "风险3", "风险4", "风险5", "合计"],
        ["价值1", "5.74%", "5.99%", "10.57%", "10.55%", "10.42%", "7.05%"],
        ["价值2", "5.46%", "8.38%", "9.69%", "13.80%", "14.34%", "9.46%"],
        ["价值3", "6.35%", "8.38%", "10.78%", "11.76%", "14.61%", "10.61%"],
        ["价值4", "8.04%", "8.15%", "14.88%", "16.80%", "20.27%", "15.91%"],
        ["价值5", "11.54%", "10.83%", "13.31%", "19.28%", "27.05%", "22.90%"],
        ["合计", "5.92%", "7.79%", "11.70%", "15.16%", "22.02%", "13.19%"],
    ]
    add_table(s, 0.56, 4.40, 11.92, d3, [1.40, 1.752, 1.752, 1.752, 1.752, 1.752, 1.76],
              0.24, center_cols={1, 2, 3, 4, 5, 6})
    box, tf = add_textbox(s, 0.56, 5.92, 11.92, 0.3)
    add_para(tf, [("注：价值1×风险5（96 个观察）、价值5×风险1（26 个观察）样本不足 100，上屏时在热力矩阵中淡化处理，不单独引用。", 7.5, False, GRAY)],
             first=True, line_spacing=1.0)
    page_concl(s, "模型适用于排序分层；OOT 弱价值概率存在约 2.15 个百分点的低估（实际 25.75% vs 预测 23.60%），用于概率阈值或定价前需先完成校准。")

    # ================= 价值模型 P3 =================
    s = prs.slides.add_slide(blank)
    page_head(s, "应用：风险 × 价值联合决策，价值扩张以风险预算为上限",
              "价值扩张以风险预算为上限，在风险约束下优化组合收益，而非单纯提高额度。")
    box, tf = add_textbox(s, 0.56, 1.30, 11.92, 0.4)
    add_para(tf, [("应用原则：", 9.5, True, INK),
                  ("额度总量受风险预算约束时，经营问题的本质是「有限额度如何分配」；离线 Lift 已证明模型排序能力，「差异化额度优于一刀切」仍需随机试验或准实验验证增量收益。", 9.5, False, BODY)],
             first=True, line_spacing=1.05)
    quads = [
        ("优先经营客群（高价值 × 低风险 · 风险1-2 × 价值1-2）",
         "样本 18,934、占比 12%，3M30 逾期率 6.30%（金额口径 3.00%）",
         "提升额度扩大价值捕获；承接 3000–5000 优质 MACC / LACC 候选池，后续筛选为循环贷扩展候选。"),
        ("谨慎经营客群（高价值 × 高风险 · 风险4-5 × 价值1-2）",
         "3M30 逾期率 12.90%",
         "有价值但风险高——不直接放大额度；通过 EWA / 超短期、小额、即时周转产品承接，控制风险暴露。"),
        ("收益验证客群（低价值 × 低风险 · 风险1-2 × 价值4-5）",
         "低风险但价值弱",
         "通过低额度 Personal Loan 验证稳定需求、复贷表现与真实贡献。"),
        ("优先风险控制客群（低价值 × 高风险 · 风险4-5 × 价值4-5）",
         "样本 77,057、占比 48.85%，3M30 逾期率 22.25%（金额口径 15.38%）",
         "限额、提价、人工审核、谨慎准入。"),
    ]
    qw, qh, qx, qy, qgapx, qgapy = 5.87, 2.42, 0.56, 1.82, 0.16, 0.12
    for i, (t, d, st) in enumerate(quads):
        x = qx + (i % 2) * (qw + qgapx)
        y = qy + (i // 2) * (qh + qgapy)
        cell = add_rect(s, x, y, qw, qh, WHITE, line=LINEC, line_w=1.5)
        add_topbar(s, x, y, qw)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.20)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.16)
        tf.margin_bottom = Inches(0.08)
        add_para(tf, [(t, 11, True, INK)], first=True, line_spacing=1.0, space_after=5)
        add_para(tf, [(d, 8.5, True, NAVY)], line_spacing=1.05, space_after=4)
        add_para(tf, [(st, 8.5, False, BODY)], line_spacing=1.05)
    page_concl(s, "额度总量受风险预算约束，经营问题的本质是「有限额度如何分配」；离线 Lift 已证明模型的排序能力，差异化额度的增量收益仍需实验验证。")

    # ================= 价值模型 P4 =================
    s = prs.slides.add_slide(blank)
    page_head(s, "画像：价值看职业现金流，风险看行为",
              "价值由收入与职业现金流决定，风险由支出、赌博与波动等行为特征决定——价值好 ≠ 风险好。")
    box, tf = add_textbox(s, 0.56, 1.28, 11.92, 0.55)
    add_para(tf, [("入模变量（17 个银行流水特征、六类）：", 9, True, INK),
                  ("收入能力（工资收入、收入均值）、现金流活跃度（借方金额均值 / 频率）、余额稳定性（余额变异系数）、负债行为（BNPL 还款、贷款还款笔数）、转账行为（转账簇占比）、支出结构（保险支出、债务客单价）；Top 3 重要度均为收入与现金流特征、合计 51%。", 9, False, BODY)],
             first=True, line_spacing=1.05)
    add_rect(s, 5.55, 1.95, 0.01, 4.4, LINEC)
    # 左列:价值轴画像
    block(s, 0.56, 1.95, 4.84, "收入与现金流（价值1/2档 vs 4/5档）", [
        [("工资收入 ", 9, True, NAVY), ("41,139", 9, True, NAVY), (" vs ", 9, False, BODY), ("11,130", 9, True, NAVY),
         ("（近4倍）；总收入 140,807 vs 59,771；工资入账 20.5 vs 10.6 次；主行业入账 40,916 vs 16,067；行业识别缺失率 3.5% vs 30.7%；Centrelink 2.1% vs 14.6%（价值5档 32.9%）；发薪周期标准差 2.72 vs 3.79。", 9, False, BODY)],
    ])
    block(s, 0.56, 3.10, 4.84, "支出、盈余与负债反证", [
        [("总支出 139,715 vs 58,603、可支配盈余 ", 9, False, BODY), ("1,334 vs 84", 9, True, NAVY),
         ("（差人群支出占收入近九成）；信贷还款 20,439 vs 5,963、29.4 vs 13.5 次——差人群信贷活动规模反而更小，不支持「价值差 = 负债更重 / 多头更多」。", 9, False, BODY)],
    ])
    block(s, 0.56, 4.30, 4.84, "家庭与渠道（结构性差异，不作因果解读）", [
        [("couple 家庭 43.8% vs 29.1%；Lead Partner 19.3% vs 11.6%、Owned Channels 11.8% vs 7.9%；差人群 Google 57.4%、Organic 21.5% 占比更高；平均年龄 37.6 vs 35.3。", 9, False, BODY)],
    ])
    # 右列:错配
    block(s, 6.28, 1.95, 6.42, "价值 ≠ 风险的证据", [
        [("工资收入随价值变差下降（41,139 → 11,130），随风险变差反而上升（风险好 ", 9, False, BODY),
         ("20,922", 9, True, NAVY), (" vs 风险差 ", 9, False, BODY), ("30,752", 9, True, NAVY),
         ("）——两个模型看的是不同的东西。", 9, False, BODY)],
    ])
    block(s, 6.28, 2.85, 6.42, "价值好但风险差（收入端强、行为端弱）", [
        [("总收入 ", 9, False, BODY), ("191,383", 9, True, NAVY),
         ("，高于「价值好且风险好」的 96,943；总支出 193,711、收支几乎打平；赌博支出占比 5.2% vs 2.5%（金额 7,416 vs 2,398）；余额异常天数 5.13 vs 3.82；收入波动 CV 2.10 vs 1.83。", 9, False, BODY)],
    ])
    block(s, 6.28, 4.10, 6.42, "价值差但风险好（低收入但克制保守）", [
        [("工资收入仅 ", 9, False, BODY), ("8,750", 9, True, NAVY),
         ("，Centrelink 占比 22.1%；总支出 43,062，明显低于「价值差且风险差」的 75,496；赌博占比 1.3%、余额波动 478；还款活动反而更活跃（还款金额 7,582 vs 4,133，次数 15.9 vs 10.8）。", 9, False, BODY)],
    ])
    # 表4 Top10(左) + 注(右)
    box, tf = add_textbox(s, 0.56, 5.55, 5.90, 0.25)
    add_para(tf, [("表4：入模变量重要度 Top10", 10.5, True, INK)], first=True, line_spacing=1.0)
    d4 = [
        ["排名", "含义", "重要度"],
        ["1", "近 182 天工资收入汇总金额", "20.25%"],
        ["2", "近 56 天收入均值", "15.60%"],
        ["3", "近 168 天借方金额均值", "15.23%"],
        ["4", "先买后付（BNPL）还款金额总和", "6.98%"],
        ["5", "近 28 天转账簇金额占比", "6.32%"],
        ["6", "个人贷款还款单机构最大累计笔数", "4.98%"],
        ["7", "近 56 天余额变异系数", "4.91%"],
        ["8", "近 168 天个人贷款放款金额均值", "4.75%"],
        ["9", "近 14 天保险出账金额", "4.56%"],
        ["10", "债务支出簇近 168 天平均客单价", "3.60%"],
    ]
    add_table(s, 0.56, 5.83, 5.90, d4, [0.70, 3.30, 1.90], 0.15, header_size=8, data_size=8,
              center_cols={0, 2})
    box, tf = add_textbox(s, 6.78, 5.83, 5.70, 0.6)
    add_para(tf, [("注：Top 3 均为收入与现金流特征（合计 51%），第 4–10 名为负债与支出类，与「价值看收入端、风险看行为端」的核心观点一致。", 8, False, GRAY)],
             first=True, line_spacing=1.05)
    page_concl(s, "两类错配人群应单独观察其真实逾期率、额度使用率与复贷表现后再定策略，不宜仅按价值或风险单边处置。")

    # ================= 银行流水 P1 =================
    s = prs.slides.add_slide(blank)
    page_head(s, "把原始流水翻译成财务画像：输入与输出",
              "银行流水是信审最重要的原料，但原料没有标签；模型把每笔交易翻译成收入与负债，支撑信审决策。")
    add_rect(s, 5.55, 1.32, 0.01, 5.3, LINEC)
    block(s, 0.56, 1.32, 4.84, "为什么重要", [
        [("澳洲 Responsible Lending 监管下，信审必须核实客户财务状况，核心公式 Serviceability = 收入 − 必要支出 − 负债还款；新客价值模型最重要的特征全部来自流水，老客风险模型里流水子模型占 ", 9, False, BODY),
         ("68.25%", 9, True, NAVY), (" 权重——流水翻译的质量，决定上层模型的天花板。", 9, False, BODY)],
    ])
    block(s, 0.56, 2.30, 4.84, "输入是什么", [
        [("一个申请名下的全部银行流水。每笔交易只有 4 个信息——时间、金额、收支方向、一段文字描述（如 PAYROLL ACME PTY LTD、BILL PAY Fair Go Finance），没有任何现成的「这是工资 / 这是贷款」标签。", 9, False, BODY)],
    ])
    block(s, 0.56, 3.35, 4.84, "输出是什么", [
        [("① 每笔交易打上分类标签 + 交易对手（共 ", 9, False, BODY), ("35 个类别", 9, True, NAVY),
         ("，如 Wages / SACC Loans / 消费）；② 收入流（月收入估算、预计下次发薪日）；③ 负债流——把外部金融机构的放款与还款交易，聚合成一条条「贷款档案」；④ 分类汇总成客户财务画像，下游据此计算可支配盈余、完成信审。", 9, False, BODY)],
    ])
    block(s, 0.56, 4.65, 4.84, "翻译示例", [
        [("PAYROLL ACME PTY LTD → ACME 公司发放的工资，归入收入流；BILL PAY Fair Go Finance → 偿还 Fair Go Finance 的 SACC 贷款，归入负债流。", 9, False, BODY)],
    ])
    block(s, 6.28, 1.32, 6.42, "输出示例——收入流（演示数据，非真实客户）", [
        [("同一雇主的连续发薪被串成一条收入流——ACME 公司每两周发薪 ", 9, False, BODY),
         ("A$2,600", 9, True, NAVY), ("，月收入估算约 ", 9, False, BODY), ("A$5,600", 9, True, NAVY),
         ("，预计下次发薪 9 月 11 日。", 9, False, BODY)],
    ])
    block(s, 6.28, 2.30, 6.42, "输出示例——负债流，核心能力（演示数据，非真实客户）", [
        [("银行流水可以还原客户在外部机构的完整多头负债档案。例如 MoneySpot 给该客户发过 3 笔贷款——A$1,200（3月，已结清）、A$1,500（6月，已结清）、A$1,500（7月，在贷）；在贷这笔按 8 期计划、每期约 A$255，推算含费成本率约 ", 9, False, BODY),
         ("36%", 9, True, NAVY), ("，属 SACC 高成本产品档。汇总到客户级：外部机构 3 家、在贷 3 笔（SACC 2 笔 + BNPL 1 笔）、近 90 天新借款 A$2,300。", 9, False, BODY)],
    ])
    block(s, 6.28, 4.45, 6.42, "指标体系的丰富性", [
        [("定价档位（从放款与还款节奏推算利率）、多头与负债压力（机构数、借款频率、剩余期数）、还款能力（月还款 / 月收入）、资金饥渴度（近 90 天新增借款）等，全部由流水客观还原，不依赖客户自报。", 9, False, BODY)],
    ])
    block(s, 6.28, 5.35, 6.42, "在链路中的位置", [
        [("IDA 发起申请 → IDP 调度 → 流水模型（BS-CAT）→ S-CALC 计算可支配盈余 → 风险模型与信审决策。", 9, False, BODY)],
    ])
    page_concl(s, "从流水还原收入与负债，把信审核实从「客户自报」升级为「流水客观还原」。")

    # ================= 银行流水 P2 =================
    s = prs.slides.add_slide(blank)
    page_head(s, "模型怎么工作、怎么运维：9 个工位的流水线",
              "架构的本质，是把专家经验变成可改、可查、可回滚的规则资产——模型越用越准，而非交付后冻结。")
    box, tf = add_textbox(s, 0.56, 1.28, 11.92, 0.5)
    add_para(tf, [("架构——9 个工位依次过一遍流水：", 9.5, True, INK),
                  ("每笔交易按顺序经过各工位，被最合适的工位认领并打上标签；收入、负债两个工位再把认领结果串成「流」。", 9.5, False, BODY)],
             first=True, line_spacing=1.05)
    nodes = ["原始流水", "① 商户识别", "② 转账", "③ 拒付", "④ 收入", "⑤ 负债",
             "⑥ 其他进账", "⑦ 费用", "⑧ 房租", "⑨ 兜底", "分类结果 +\n收入/负债流"]
    strip_y, strip_h = 1.86, 0.70
    add_rect(s, 0.56, strip_y, 11.92, strip_h, LGRAY)
    nw, nh, gap = 0.98, 0.46, 0.10
    total_w = len(nodes) * nw + (len(nodes) - 1) * gap
    x0 = 0.56 + (11.92 - total_w) / 2
    ny = strip_y + (strip_h - nh) / 2
    for i, name in enumerate(nodes):
        x = x0 + i * (nw + gap)
        first_last = (i == 0 or i == len(nodes) - 1)
        node = add_rect(s, x, ny, nw, nh, NAVY if first_last else WHITE,
                        line=None if first_last else NAVY, line_w=1.0)
        tf = node.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)
        for j, ln in enumerate(name.split('\n')):
            add_para(tf, [(ln, 8.5, True, WHITE if first_last else NAVY)],
                     first=(j == 0), align=PP_ALIGN.CENTER, line_spacing=1.0)
        if i < len(nodes) - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + nw + 0.012),
                                    Inches(ny + nh / 2 - 0.045), Inches(gap - 0.024), Inches(0.09))
            ar.fill.solid()
            ar.fill.fore_color.rgb = GRAY
            ar.line.fill.background()
            ar.shadow.inherit = False
    add_rect(s, 5.55, 2.80, 0.01, 2.7, LINEC)
    block(s, 0.56, 2.80, 4.84, "设计思考①：为什么用「规则 + 知识库」，而不是训练一个算法模型", [
        [("信贷场景的刚需是可解释、可审计：每一笔交易为什么这样分类，都能指到具体规则，被拒客户申诉与监管检查都需要；变更成本极低：改一条规则就是改一行 Excel，不用重新训练、重新部署；能力随业务累积：商户库已滚到 ", 8.5, False, BODY),
         ("87.6 万家", 8.5, True, NAVY), ("，模型随使用越滚越强。", 8.5, False, BODY)],
    ])
    block(s, 0.56, 4.35, 4.84, "设计思考②：为什么拆成 9 个工位", [
        [("职责单一，问题能定位到具体工位；互不干扰：改负债规则不会误伤收入识别；类别冲突靠「顺序 + 保护规则」解决：收入是信审最重要的信号，规则上保护它不被其他工位覆盖；每个工位独立测试、独立迭代，回归检查精确到工位层。", 8.5, False, BODY)],
    ])
    block(s, 6.28, 2.80, 6.42, "设计思考③：为什么做四层基线回归", [
        [("「最终结果没变」不等于「没影响」：某工位改了规则但没当上最终赢家，最终结果看不出变化，隐患会在未来爆发；四层回归（最终结果 / 各工位认领 / 配置指纹 / 汇总指标）让任何调整都可追溯、可回滚；这也是敢高频更新知识库的前提——更新便宜且安全。", 8.5, False, BODY)],
    ])
    block(s, 6.28, 4.35, 6.42, "设计思考④：知识库就是这些 Excel", [
        [("全部规则外置在 24 张规则表中，业务同事可以直接打开修改——最大的「商户库」已收录 ", 8.5, False, BODY),
         ("87.6 万家", 8.5, True, NAVY),
         ("澳洲商户，每行是「商户名 + 关键词变体 + 分类」；收入表定义工资识别模式；负债表收录贷款机构与产品类型；房租规则带置信度评分。改表即改规则，不用动代码。", 8.5, False, BODY)],
    ])
    block(s, 0.56, 5.65, 5.84, "怎么运维：「2 人 + 1 个 AI 助手」", [
        [("模型团队同学负责主动运维（质检监控、分类漂移与模型波动），并处理人审团队以 ticket 反馈的 case，给出解决方案与代码修改；AI 助手负责两类监测：分类异常波动与大量缺失的主动提示、市面上新商户的新增监测。", 8.5, False, BODY)],
    ])
    block(s, 6.64, 5.65, 5.84, "相对 illion 的收益", [
        [("逻辑透明、问题可定位；粒度可控（SACC / Non-SACC / BNPL / LOC / EWA 可拆分）；口径可配置；Open Banking 等新数据源可直接接入。", 8.5, False, BODY)],
    ])
    page_concl(s, "架构的每一项取舍都指向同一件事——让规则资产可持续演进：改得动、看得清、回得去。")

    prs.save(PPTX)
    print("saved:", PPTX, "slides:", len(prs.slides))


if __name__ == "__main__":
    main()
