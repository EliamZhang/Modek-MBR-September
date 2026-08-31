# -*- coding: utf-8 -*-
"""把月会PPT 第4页重建为图片版「四类客群矩阵」:
- 页顶改为「四类客群矩阵」金色竖条标题(去掉原页面大标题/引导语/结论条)
- 四象限客群卡,每卡含: 顶部两标签胶囊 / 客群名 / 对应产品框 / 描述 / 指标卡 / 底部标签行
- 四色主题: 优先经营(绿) 谨慎经营(金) 收益验证(蓝) 风险控制(红)
- 数据与正式版MD P3 完全一致,未改动。
参考 sync_profile_pages.py 的 helper 与配色。"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout.reconfigure(encoding='utf-8')

PPTX = r"c:\Users\zhangyuliang02\Desktop\临时文件\9月MBR\月会PPT.pptx"

INK = RGBColor(0x0A, 0x0A, 0x0A)
BODY = RGBColor(0x3D, 0x3C, 0x3A)
GRAY = RGBColor(0x60, 0x60, 0x60)
NAVY = RGBColor(0x1F, 0x4E, 0x79)
GOLD = RGBColor(0xD6, 0xA0, 0x00)
LGOLD = RGBColor(0xFA, 0xF3, 0xE0)
LINEC = RGBColor(0xD6, 0xD6, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"


# ---------- 基础 helper ----------
def set_run(run, size, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_para(tf, runs, first=False, align=PP_ALIGN.LEFT, space_after=0, line_spacing=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        set_run(r, size, bold, color)
    return p


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def add_rect(slide, x, y, w, h, fill, line=None, line_w=1.5, rounded=False, radius=0.03):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def pill(slide, x, y, w, h, text, fill, line, tcolor, size=9, bold=True):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    p.adjustments[0] = 0.5
    p.fill.solid()
    p.fill.fore_color.rgb = fill
    if line is None:
        p.line.fill.background()
    else:
        p.line.color.rgb = line
        p.line.width = Pt(1)
    p.shadow.inherit = False
    tf = p.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    par = tf.paragraphs[0]
    par.alignment = PP_ALIGN.CENTER
    par.line_spacing = 1.0
    r = par.add_run()
    r.text = text
    set_run(r, size, bold, tcolor)
    return p


def cw_est(ch, pt):
    if ord(ch) > 0x2E7F:
        return pt / 72.0
    if 0x20 <= ord(ch) <= 0x7E:
        return pt / 72.0 * 0.52
    return pt / 72.0 * 0.9


def text_w(text, pt):
    return sum(cw_est(c, pt) for c in text)


def clear_slide(slide):
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)


# ---------- 卡片内容 ----------
QUADS = [
    dict(
        name="优先经营客群", accent=RGBColor(0x3E, 0x7A, 0x55), light=RGBColor(0xEE, 0xF6, 0xF0),
        risk="风险 1-2 × 价值 1-2", val="价值好 × 风险好",
        product_label="对应产品", product="3000-5000 优质 MACC / LACC 候选池",
        desc="风险低、价值好，适合作为重点经营客群。当前优先承接 3000-5000 优质 MACC，后续可从中筛选更优质、更稳定人群，作为 LACC / 循环贷的扩展候选。",
        m1=[("样本数", "18,934", "num"), ("样本占比", "12.00%", "num"), ("3M30人数逾期率", "6.30%", "num")],
        m2=[("3M30金额逾期率", "3.00%", "num")],
        tags=["优先经营", "优质 MACC", "LACC 候选", "循环贷候选"],
    ),
    dict(
        name="谨慎经营客群", accent=RGBColor(0xB8, 0x86, 0x0B), light=RGBColor(0xFB, 0xF4, 0xE2),
        risk="风险 4-5 × 价值 1-2", val="价值好 × 风险差",
        product_label="对应产品", product="EWA / 超短期产品",
        desc="客户具备一定经营价值，但风险压力已经抬升。不适合直接放款大额度或期限，更适合通过短期、小额、即时周转产品承接需求，并控制风险暴露。",
        m1=[("定位", "有价值但风险高", "txt"), ("3M30人数逾期率", "12.90%", "num")],
        m2=[],
        tags=["短期承接", "小额控制", "即时周转", "风险暴露控制"],
    ),
    dict(
        name="收益验证客群", accent=RGBColor(0x1F, 0x4E, 0x79), light=RGBColor(0xEF, 0xF4, 0xFA),
        risk="风险 1-2 × 价值 4-5", val="价值差 × 风险好",
        product_label="对应产品", product="低额度 Personal Loan",
        desc="风险相对可控，但收入、额度承接和转化能力偏弱。不适合作为强经营客群，可通过低额度 Personal Loan 验证稳定需求、复贷表现和真实贡献。",
        m1=[("定位", "低风险但价值弱", "txt"), ("观察重点", "稳定需求 / 复贷 / 真实贡献", "txt")],
        m2=[],
        tags=["低额度验证", "收益观察", "复贷跟踪", "贡献评估"],
    ),
    dict(
        name="优先风险控制客群", accent=RGBColor(0xA9, 0x44, 0x34), light=RGBColor(0xF9, 0xEC, 0xE9),
        risk="风险 4-5 × 价值 4-5", val="价值差 × 风险差",
        product_label="对应产品", product="原则上不重点承接",
        desc="同时处在高风险、弱价值位置，风险暴露明显。该客群不建议绑定扩展型产品，重点应放在限额、提价、人工审核和谨慎准入。",
        m1=[("样本数", "77,057", "num"), ("样本占比", "48.85%", "num"), ("3M30人数逾期率", "22.25%", "num")],
        m2=[("3M30金额逾期率", "15.38%", "num")],
        tags=["限额", "提价", "人工审核", "谨慎准入"],
    ),
]


def metric(slide, x, y, w, h, label, value, kind, accent, light):
    box = add_rect(slide, x, y, w, h, light, line=None)
    add_rect(slide, x, y, w, 0.045, accent)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.09)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.02)
    add_para(tf, [(label, 8, False, GRAY)], first=True, line_spacing=1.0, space_after=1)
    vs = 11 if kind == "num" else 9
    add_para(tf, [(value, vs, True, accent)], line_spacing=1.0)


def build_card(slide, cx, cy, W, H, q):
    pl, pr = 0.18, 0.16
    cw = W - pl - pr
    # 卡底
    add_rect(slide, cx, cy, W, H, q["light"], line=q["accent"], line_w=1.5,
             rounded=True, radius=0.035)
    # 顶部两个标签胶囊
    tx = cx + pl
    ty = cy + 0.14
    th = 0.26
    for text, is_risk in ((q["risk"], True), (q["val"], False)):
        w = text_w(text, 9) + 0.30
        fill = q["light"] if is_risk else q["accent"]
        line = q["accent"] if is_risk else None
        tcol = q["accent"] if is_risk else WHITE
        pill(slide, tx, ty, w, th, text, fill, line, tcol, size=9, bold=True)
        tx += w + 0.10
    # 客群名
    add_textbox(slide, cx + pl, cy + 0.46, cw, 0.30)
    box, tf = add_textbox(slide, cx + pl, cy + 0.46, cw, 0.32)
    add_para(tf, [(q["name"], 14, True, q["accent"])], first=True, line_spacing=1.0)
    # 产品框
    add_rect(slide, cx + pl, cy + 0.82, cw, 0.34, WHITE, line=q["accent"],
             line_w=1.0, rounded=True, radius=0.10)
    box, tf = add_textbox(slide, cx + pl + 0.12, cy + 0.85, cw - 0.24, 0.28)
    add_para(tf, [(q["product_label"] + "  ", 8.5, False, GRAY),
                  (q["product"], 9.5, True, q["accent"])], first=True, line_spacing=1.0)
    # 描述
    box, tf = add_textbox(slide, cx + pl, cy + 1.24, cw, 0.70)
    add_para(tf, [(q["desc"], 9, False, BODY)], first=True, line_spacing=1.15)
    # 指标区
    my = cy + 2.06
    mh = 0.40
    gap = 0.06
    rows = [q["m1"], q["m2"]]
    for row in rows:
        if not row:
            continue
        n = len(row)
        w_each = (cw - (n - 1) * gap) / n
        x = cx + pl
        for label, value, kind in row:
            metric(slide, x, my, w_each, mh, label, value, kind,
                   q["accent"], RGBColor(0xFF, 0xFF, 0xFF))
            x += w_each + gap
        my += mh + 0.06
    # 底部标签行
    tx = cx + pl
    ty = cy + H - 0.26
    th = 0.24
    for tag in q["tags"]:
        w = text_w(tag, 8) + 0.26
        pill(slide, tx, ty, w, th, tag, q["accent"], None, WHITE, size=8, bold=True)
        tx += w + 0.08


def main():
    prs = Presentation(PPTX)
    s = prs.slides[3]  # 第4页
    clear_slide(s)

    # 页顶标题: 金竖条 + 四类客群矩阵
    add_rect(s, 0.56, 0.28, 0.08, 0.24, GOLD)
    box, tf = add_textbox(s, 0.70, 0.26, 8.00, 0.30)
    add_para(tf, [("四类客群矩阵", 16, True, NAVY)], first=True, line_spacing=1.0)

    W, H = 5.90, 3.24
    x0, x1 = 0.56, 6.66
    y0, y1 = 0.72, 4.08
    build_card(s, x0, y0, W, H, QUADS[0])  # 优先经营
    build_card(s, x1, y0, W, H, QUADS[1])  # 谨慎经营
    build_card(s, x0, y1, W, H, QUADS[2])  # 收益验证
    build_card(s, x1, y1, W, H, QUADS[3])  # 风险控制

    prs.save(PPTX)
    print("saved 第4页 → 四类客群矩阵")


if __name__ == "__main__":
    main()
