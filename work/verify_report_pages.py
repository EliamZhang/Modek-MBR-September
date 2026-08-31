# -*- coding: utf-8 -*-
"""对新增 6 页做结构化布局验证:折行估算、文本溢出、边界、重叠。"""
from pptx import Presentation

PPTX = r"c:\Users\zhangyuliang02\Desktop\临时文件\9月MBR\Model MBR August.pptx"


def cw(ch, pt):
    if ord(ch) > 0x2E7F:
        return pt / 72.0
    if 0x20 <= ord(ch) <= 0x7E:
        return pt / 72.0 * 0.52
    return pt / 72.0 * 0.9


def est_lines(text, pt, avail):
    if not text.strip():
        return 0
    cur = 0.0
    n = 1
    for ch in text:
        cur += cw(ch, pt)
        if cur > avail:
            n += 1
            cur = cw(ch, pt)
    return n


prs = Presentation(PPTX)
SW, SH = 13.33, 7.5
for si, s in enumerate(prs.slides, 1):
    if si == 1:
        continue
    print("=" * 70)
    print("SLIDE", si)
    shapes = list(s.shapes)
    issues = 0
    for sh in shapes:
        l = sh.left / 914400
        t = sh.top / 914400
        w = sh.width / 914400
        h = sh.height / 914400
        r = l + w
        b = t + h
        if l < -0.02 or t < -0.02 or r > SW + 0.02 or b > SH + 0.02:
            print("  ** OUT OF BOUNDS: (%.2f,%.2f)-(%.2f,%.2f)" % (l, t, r, b))
            issues += 1
        if sh.has_table:
            tbl = sh.table
            rh = sum(tbl.rows[i].height for i in range(len(tbl.rows))) / 914400
            if rh > h + 0.01:
                print("  ** TABLE taller than frame: rows %.2f > frame %.2f" % (rh, h))
                issues += 1
            continue
        if not (sh.has_text_frame and sh.text_frame.text.strip()):
            continue
        # 文本区可用宽 = 形状宽 - 左右边距
        ml = sh.text_frame.margin_left / 914400 if sh.text_frame.margin_left else 0
        mr = sh.text_frame.margin_right / 914400 if sh.text_frame.margin_right else 0
        mt = sh.text_frame.margin_top / 914400 if sh.text_frame.margin_top else 0
        mb = sh.text_frame.margin_bottom / 914400 if sh.text_frame.margin_bottom else 0
        avail_w = w - ml - mr
        avail_h = h - mt - mb
        total_h = 0.0
        for para in sh.text_frame.paragraphs:
            ptxt = para.text
            if not ptxt.strip():
                continue
            ps = 9
            if para.runs and para.runs[0].font.size:
                ps = para.runs[0].font.size.pt
            n = est_lines(ptxt, ps, avail_w)
            total_h += n * (ps / 72.0) * 1.15
            if para.space_after is not None:
                total_h += para.space_after.pt / 72.0
        if total_h > avail_h + 0.06:
            txt = sh.text_frame.text.replace(chr(10), ' | ')[:38]
            print("  ** OVERFLOW: text %.2f > avail %.2f :: %s" % (total_h, avail_h, txt))
            issues += 1
    # 重叠:含文本形状两两
    text_shapes = [sh for sh in shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    for i in range(len(text_shapes)):
        for j in range(i + 1, len(text_shapes)):
            a, b = text_shapes[i], text_shapes[j]
            ox = min(a.left + a.width, b.left + b.width) - max(a.left, b.left)
            oy = min(a.top + a.height, b.top + b.height) - max(a.top, b.top)
            if ox > 91440 * 0.05 and oy > 91440 * 0.05:
                print("  ** OVERLAP: [%s] x [%s]" % (
                    a.text_frame.text.replace(chr(10), '|')[:24],
                    b.text_frame.text.replace(chr(10), '|')[:24]))
                issues += 1
    print("  issues:", issues)
