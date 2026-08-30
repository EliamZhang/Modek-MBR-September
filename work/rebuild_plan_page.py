# -*- coding: utf-8 -*-
"""重建 Model MBR August.pptx 单页总览:四项目标 + 四大领域 2×2(深蓝+金)。
本版:压缩文字、条目在块内垂直均匀分布,缓解 4 块密度不均。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

PPTX = r"c:\Users\zhangyuliang02\Desktop\临时文件\9月MBR\Model MBR August.pptx"

INK = RGBColor(0x0A, 0x0A, 0x0A)
BODY = RGBColor(0x3D, 0x3C, 0x3A)
GRAY = RGBColor(0x60, 0x60, 0x60)
NAVY = RGBColor(0x1F, 0x4E, 0x79)
GOLD = RGBColor(0xD6, 0xA0, 0x00)
LBLUE = RGBColor(0xE3, 0xEB, 0xF4)
LGOLD = RGBColor(0xFA, 0xF3, 0xE0)
LINEC = RGBColor(0xD6, 0xD6, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"


def set_run(run, size, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_para(tf, runs, first=False, align=PP_ALIGN.LEFT, space_after=0,
             line_spacing=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    for text, size, bold, color in runs:
        r = p.add_run()
        r.text = text
        set_run(r, size, bold, color)
    return p


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def add_rect(slide, x, y, w, h, fill, line=None, line_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_topbar(slide, x, y, w):
    return add_rect(slide, x, y, w, 0.06, NAVY)


def add_section_title(slide, y, text):
    add_rect(slide, 0.56, y + 0.02, 0.07, 0.22, GOLD)
    box, tf = add_textbox(slide, 0.70, y, 9.0, 0.30)
    add_para(tf, [(text, 12, True, INK)], first=True, line_spacing=1.0)


def main():
    prs = Presentation(PPTX)
    slide = prs.slides[0]
    spTree = slide.shapes._spTree
    for child in list(spTree):
        if child.tag != qn('p:nvGrpSpPr'):
            spTree.remove(child)

    # ---- 标题 / 引导语 ----
    box, tf = add_textbox(slide, 0.66, 0.23, 11.92, 0.40)
    add_para(tf, [("体系规划:四项目标,四大领域", 20, True, INK)],
             first=True, line_spacing=1.3)
    box, tf = add_textbox(slide, 0.56, 0.64, 11.92, 0.48)
    add_para(tf, [("规模、结构、质量、效率四项目标并重;数据基建、投放、运营、风险四大领域共享同一分层体系。", 18, False, GRAY)],
             first=True, line_spacing=1.4)
    add_rect(slide, 0.56, 1.14, 11.92, 0.012, LINEC)

    # ---- 四项目标 ----
    add_section_title(slide, 1.28, "四项目标:规模 · 结构 · 质量 · 效率")
    goals = [
        ("规模", "支撑放款增长", "拓展 APP、新渠道与复贷"),
        ("结构", "NON SAAC 全年 ≥ 40%", "非 SACC 产品占比目标"),
        ("质量", "风险表现控制在 7.1%", "整体风险水平控制"),
        ("效率", "提升转化率与投放 ROI", "以 LTV 统筹获客、额度、定价、权益"),
    ]
    gw, gh, gap, gy = 2.85, 1.08, 0.19, 1.60
    for i, (name, metric, desc) in enumerate(goals):
        x = 0.56 + i * (gw + gap)
        card = add_rect(slide, x, gy, gw, gh, LBLUE)
        add_topbar(slide, x, gy, gw)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.16)
        tf.margin_bottom = Inches(0.08)
        add_para(tf, [(name, 10.5, True, INK)], first=True, line_spacing=1.0, space_after=4)
        add_para(tf, [(metric, 13, True, NAVY)], line_spacing=1.0, space_after=4)
        add_para(tf, [(desc, 8.5, False, BODY)], line_spacing=1.0)

    # ---- 四大领域:2×2 网格,条目垂直均匀分布 ----
    add_section_title(slide, 2.90, "四大领域:共享同一分层体系")
    bw, bh, bx, by1, gapx, gapy = 5.87, 1.76, 0.56, 3.20, 0.16, 0.10
    LINES_GAP = 0.175  # 8.5pt 条目行基线高(含行距)

    def field_card(x, y, title, lines):
        card = add_rect(slide, x, y, bw, bh, WHITE, line=LINEC, line_w=1.5)
        add_topbar(slide, x, y, bw)
        add_rect(slide, x + 0.16, y + 0.16, 0.06, 0.20, GOLD)
        # 可用条目区高度 = 块高 - 标题区(标题+上边距) - 下边距
        title_h = 0.32  # 12pt 标题行 + space_after 5pt
        avail = bh - 0.14 - title_h - 0.08
        n_units = len(lines)
        row_h = 8.5 / 72.0 * 1.0      # 条目行高(8.5pt 单行)
        # 垂直两端对齐:每行行距 = 可用高减去所有行高后均分(上限适中,不松散也不过挤)
        gap = (avail - n_units * row_h) / max(n_units - 1, 1)
        gap = min(max(gap, 0.0), 0.12)  # 上限 0.12",留白主要落在下部
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.28)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.14)
        tf.margin_bottom = Inches(0.08)
        add_para(tf, [(title, 12, True, INK)], first=True, line_spacing=1.0, space_after=5)
        for ln in lines:
            if len(ln) == 1:  # 组标题
                add_para(tf, [(ln[0], 9.5, True, NAVY)], line_spacing=1.0, space_after=gap)
            else:
                text, star = ln
                if star:
                    add_para(tf, [("★ ", 8.5, True, GOLD), (text, 8.5, True, NAVY)],
                             line_spacing=1.0, space_after=gap)
                else:
                    add_para(tf, [(text, 8.5, False, BODY)], line_spacing=1.0, space_after=gap)

    bx2 = bx + bw + gapx
    by2 = by1 + bh + gapy

    # 左上 · 数据基建(精简,2组标题+5条)
    field_card(bx, by1, "数据基建", [
        ("流水识别与统一底座",),
        ("BS-CAT 生产化与并行验证", True),
        ("打通 IDP2 全链路:IDA → IDP → BS-CAT → S-CALC → 落库 → 人工复核", True),
        ("商户 / 类别知识库智能体", False),
        ("归因基建",),
        ("投放归因与回传:渠道/Campaign overlap、Web 追踪、回传质量监控", True),
        ("触达归因与效果监控:短信/邮件/Push/APP 曝光 → 放款", True),
    ])
    # 右上 · 投放市场
    field_card(bx2, by1, "投放市场", [
        ("人群圈选:风险 × 意愿 × 价值分层输出人群包", False),
        ("Lead Partner 线索评估", False),
        ("边际 CPS 与扩量决策", True),
        ("Campaign 早期质量预测", False),
    ])
    # 左下 · 用户运营(2 组)
    field_card(bx, by2, "用户运营", [
        ("客户价值与定价",),
        ("盈利性评估", False), ("自动寻优定价", False), ("Fundo Score", False),
        ("触达与智能经营",),
        ("触达策略与人群编排", True), ("运营监控 AI 智能体", False),
        ("结清客户复贷召回", True), ("沉睡客户识别唤醒", False),
        ("竞品经营洞察 AI 智能体", False),
    ])
    # 右下 · 风险管理(压缩到单行)
    field_card(bx2, by2, "风险管理", [
        ("新老客 PD 模型迭代:增强头尾区分,平衡通过率与风险", True),
        ("风险 × 价值联合额度策略:交叉分箱自动寻优,输出最优额度与承接标签", True),
    ])

    # ---- 结论条 ----
    cy, ch = 6.82, 0.42
    concl = add_rect(slide, 0.56, cy, 11.92, ch, LGOLD)
    add_rect(slide, 0.56, cy, 0.07, ch, GOLD)
    tf = concl.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.24)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    add_para(tf, [("四项目标统领全局——数据基建打底,投放、运营、风险四大领域在统一数据与分层体系上协同建设。", 9, True, INK)],
             first=True, line_spacing=1.1)

    prs.save(PPTX)
    print("saved:", PPTX)


if __name__ == "__main__":
    main()
