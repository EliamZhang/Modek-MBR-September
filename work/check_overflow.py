# -*- coding: utf-8 -*-
"""按真实字体字宽估算四大领域块每段行数,判定是否溢出。"""
from pptx import Presentation

p = Presentation(r"c:\Users\zhangyuliang02\Desktop\临时文件\9月MBR\Model MBR August.pptx")
s = p.slides[0]


def cw(ch, pt):
    # 全角(中文/全角标点)≈ pt/72 in;半角(字母/数字/半角标点)≈ 0.52 倍
    if ord(ch) > 0x2E7F:
        return pt / 72.0
    if 0x20 <= ord(ch) <= 0x7E:
        return pt / 72.0 * 0.52
    return pt / 72.0 * 0.9


def est_lines(text, pt, avail):
    if not text.strip():
        return 0
    cur = 0.0
    n = 1
    for ch in text:
        cur += cw(ch, pt)
        if cur > avail:
            n += 1
            cur = cw(ch, pt)
    return n


for sh in s.shapes:
    if sh.has_text_frame:
        t = sh.text_frame.text
        name = next((n for n in ['数据基建', '投放市场', '用户运营', '风险管理'] if t.startswith(n)), None)
        if name:
            cardw = sh.width / 914400
            avail = cardw - 0.28 - 0.14   # 左右内边距
            bh = sh.height / 914400
            availh = bh - 0.14 - 0.08     # 上下内边距
            totalh = 0.0
            lines = 0
            for para in sh.text_frame.paragraphs:
                ptxt = para.text
                if not ptxt.strip():
                    continue
                ps = para.runs[0].font.size.pt if para.runs and para.runs[0].font.size else 8.5
                n = est_lines(ptxt, ps, avail)
                lines += n
                totalh += n * (ps / 72.0)
                if para.space_after:
                    totalh += para.space_after.pt / 72.0
            delta = totalh - availh
            st = "OK" if totalh <= availh + 0.05 else "OVERFLOW %.2f in" % delta
            print("%s: 宽%.2f 可用高%.2f | 行数%d 文本高%.2f -> %s" % (name, cardw, availh, lines, totalh, st))
