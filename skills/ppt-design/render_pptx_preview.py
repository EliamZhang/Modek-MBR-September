# -*- coding: utf-8 -*-
"""渲染 PPTX 指定页为近似预览 PNG(供人工查看布局,无需 PowerPoint)。

用途:PowerPoint COM 不可用(或需要快速预览)时,用 python-pptx 读取真实形状
坐标 + run 颜色,按字宽估测折行,用 PIL 重绘到 PNG。是「近似」预览——位置、
颜色、字号来自真实数据,但字体渲染与 PowerPoint 有差异,只用于判断布局、
是否溢出、疏密是否均匀,不用于像素级校验。

用法:
    python render_pptx_preview.py <pptx路径> <输出png路径> [页序号 1..n] [图宽px]

依赖:pip install python-pptx pillow
"""
import sys
from pptx import Presentation
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

DEFAULT_FONT = r'C:/Windows/Fonts/msyh.ttc'


def px_scale(pt):
    return max(int(pt * 2.5), 8)


def get_font(pt, path):
    try:
        return ImageFont.truetype(path, px_scale(pt))
    except Exception:
        return ImageFont.load_default()


def char_w(ch, pt):
    # 全角(中文/全角标点)≈ pt/72 in;半角(字母/数字/半角标点)≈ 0.52 倍
    if ord(ch) > 0x2E7F:
        return pt / 72.0
    if 0x20 <= ord(ch) <= 0x7E:
        return pt / 72.0 * 0.52
    return pt / 72.0 * 0.9


def wrap_lines(text, pt, avail_in):
    if not text.strip():
        return []
    cur = ''
    out = []
    for ch in text:
        if sum(char_w(c, pt) for c in cur + ch) > avail_in and cur:
            out.append(cur)
            cur = ch
        else:
            cur += ch
    if cur:
        out.append(cur)
    return out


def hex_rgb(h):
    if not h:
        return None
    h = h.lstrip('#')
    while len(h) < 6:
        h += '0'
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def shape_style(sh):
    """返回 (fill_rgb, line_rgb)。"""
    spPr = sh._element.find(qn('p:spPr'))
    fill = line = None
    if spPr is not None:
        sf = spPr.find(qn('a:solidFill'))
        if sf is not None:
            srgb = sf.find(qn('a:srgbClr'))
            if srgb is not None:
                fill = hex_rgb(srgb.get('val'))
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            nf = ln.find(qn('a:solidFill'))
            if nf is not None:
                srgb = nf.find(qn('a:srgbClr'))
                if srgb is not None:
                    line = hex_rgb(srgb.get('val'))
    return fill, line


def run_color(run):
    """run 的实际颜色(RGBColor 的 str 是 16 进制串如 '1F4E79')。"""
    try:
        if run.font.color and run.font.color.rgb is not None:
            h = str(run.font.color.rgb)
            return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        pass
    return (10, 10, 10)


def main():
    if len(sys.argv) < 3:
        print("usage: python render_pptx_preview.py <pptx> <out.png> [page 1..n] [widthpx]")
        sys.exit(1)
    pptx_path = sys.argv[1]
    out_png = sys.argv[2]
    page_no = int(sys.argv[3]) if len(sys.argv) > 3 else 1
    width_px = int(sys.argv[4]) if len(sys.argv) > 4 else 2400

    prs = Presentation(pptx_path)
    if page_no < 1 or page_no > len(prs.slides):
        print("page out of range:", page_no, "slides:", len(prs.slides))
        sys.exit(1)
    s = prs.slides[page_no - 1]
    sw, sh = prs.slide_width / 914400, prs.slide_height / 914400
    W = width_px
    H = int(W * sh / sw)
    S = W / sw
    img = Image.new('RGB', (W, H), 'white')
    d = ImageDraw.Draw(img)

    def P(v):
        return v * S

    # 先画无文本的装饰形状(顶条/竖条/分隔线/纯色块)
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            continue
        fill, line = shape_style(sh)
        l, t = sh.left / 914400, sh.top / 914400
        w, h = sh.width / 914400, sh.height / 914400
        d.rectangle([P(l), P(t), P(l + w), P(t + h)], fill=fill or (255, 255, 255),
                    outline=line)

    # 再画含文本的卡(底色 + 折行文字)
    for sh in s.shapes:
        # 表格(graphicFrame)单独绘制:网格 + 单元格文字
        if sh.has_table:
            tbl = sh.table
            l, t = sh.left / 914400, sh.top / 914400
            w, h = sh.width / 914400, sh.height / 914400
            d.rectangle([P(l), P(t), P(l + w), P(t + h)], fill=(255, 255, 255),
                        outline=(10, 10, 10), width=2)
            nrows = len(tbl.rows)
            ncols = len(tbl.columns)
            row_tops = []
            acc = t
            for ri in range(nrows):
                row_tops.append(acc)
                acc += tbl.rows[ri].height / 914400
            col_lefts = []
            acc = l
            for ci in range(ncols):
                col_lefts.append(acc)
                acc += tbl.columns[ci].width / 914400
            for ri in range(nrows):
                y0 = row_tops[ri]
                y1 = row_tops[ri + 1] if ri + 1 < nrows else t + h
                for ci in range(ncols):
                    x0 = col_lefts[ci]
                    x1 = col_lefts[ci + 1] if ci + 1 < ncols else l + w
                    cell = tbl.cell(ri, ci)
                    try:
                        cfill = cell.fill.fore_color.rgb
                        cell_bg = tuple(int(str(cfill)[i:i + 2], 16) for i in (0, 2, 4)) if cfill else None
                    except Exception:
                        cell_bg = None
                    if cell_bg and cell_bg != (255, 255, 255):
                        d.rectangle([P(x0), P(y0), P(x1), P(y1)], fill=cell_bg)
                    txt = cell.text_frame.text
                    if txt.strip():
                        d.text((P(x0 + 0.04), P(y0 + 0.01)), txt.replace('\n', ' ')[:12],
                               font=get_font(7.5, DEFAULT_FONT), fill=(61, 60, 58))
                    # 行分隔线
                    if ci == ncols - 1:
                        d.line([(P(l), P(y1)), (P(l + w), P(y1))],
                               fill=(10, 10, 10) if ri == 0 else (214, 214, 214),
                               width=3 if ri == 0 else 1)
                # 列分隔线仅表头以下细灰线
            continue
        if not (sh.has_text_frame and sh.text_frame.text.strip()):
            continue
        fill, line = shape_style(sh)
        l, t = sh.left / 914400, sh.top / 914400
        w, h = sh.width / 914400, sh.height / 914400
        d.rectangle([P(l), P(t), P(l + w), P(t + h)], fill=fill or (255, 255, 255),
                    outline=line)
        avail = w - 0.28
        cy = t + 0.12
        for para in sh.text_frame.paragraphs:
            ptxt = para.text
            if not ptxt.strip():
                continue
            # 段落字号(取第一个 run)
            ps = 8.5
            col = (10, 10, 10)
            runs = para.runs
            if runs and runs[0].font.size:
                ps = runs[0].font.size.pt
            if runs and runs[0].font.color and runs[0].font.color.rgb is not None:
                col = run_color(runs[0])
            f = get_font(ps, DEFAULT_FONT)
            for ln in wrap_lines(ptxt, ps, avail):
                d.text((P(l + 0.18), P(cy)), ln, font=f, fill=col)
                cy += ps / 72.0 * 1.15
            if para.space_after is not None and para.space_after.pt:
                cy += para.space_after.pt / 72.0

    img.save(out_png)
    print("saved:", out_png, "page", page_no)


if __name__ == "__main__":
    main()
